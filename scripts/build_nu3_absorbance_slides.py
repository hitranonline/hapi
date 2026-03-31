"""
Build a PowerPoint presentation from nu3 absorbance progression PNGs.

Each slide contains one figure with temperature annotation and brief analysis.
Uses the template from 38.34_laser_induced_plasma/template.pptx.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN


TEMPLATE_PATH = Path("/Volumes/T7/38_Presentation&Drafts/38.34_laser_induced_plasma/template.pptx")
ARTIFACTS_DIR = Path("/Volumes/T7/GitHub/hapi/artifacts")
OUTPUT_PATH = ARTIFACTS_DIR / "nu3_absorbance_temperature_slides.pptx"

PROGRESSIONS = ["nu3_0_to_1", "nu3_1_to_2", "nu3_2_to_3", "nu3_3_to_4"]
PROGRESSION_LABELS = {
    "nu3_0_to_1": "ν₃: 0→1 (fundamental)",
    "nu3_1_to_2": "ν₃: 1→2 (1st hot band)",
    "nu3_2_to_3": "ν₃: 2→3 (2nd hot band)",
    "nu3_3_to_4": "ν₃: 3→4 (3rd hot band)",
}

# (folder_name, display_label, case_description, per-progression analysis)
CASES = [
    (
        "combined_exomol_i1_absorbance_T300K_P3Torr_x0p008_L100cm_step0p001",
        "Combined ExoMol I1 + HITRAN — 300 K",
        "T = 300 K, P = 3 Torr, x = 0.008, L = 100 cm, Δν = 0.001 cm⁻¹",
        {
            "nu3_0_to_1": "Strong fundamental band. At 300 K, nearly all molecules are in v₃ = 0. "
                          "This is the only progression with significant absorbance.",
            "nu3_1_to_2": "Peak absorbance ~10⁻⁷. The Boltzmann population of v₃ = 1 at 300 K is "
                          "only 4.2×10⁻⁷ — less than 1 in a million molecules.",
            "nu3_2_to_3": "Peak absorbance ~7×10⁻¹³. The v₃ = 2 population is the square of the "
                          "v₃ = 1 factor, making this band essentially invisible.",
            "nu3_3_to_4": "Peak absorbance = 0 (peak_source = none). The v₃ = 3 population at 300 K "
                          "is below floating-point significance.",
        },
    ),
    (
        "combined_exomol_i1_absorbance_T600K_P3Torr_x0p008_L100cm_step0p001",
        "Combined ExoMol I1 + HITRAN — 600 K",
        "T = 600 K, P = 3 Torr, x = 0.008, L = 100 cm, Δν = 0.001 cm⁻¹",
        {
            "nu3_0_to_1": "Fundamental weakens compared to 300 K due to partition function dilution. "
                          "Peak drops from ~1.6 to ~0.29 as population spreads to excited states.",
            "nu3_1_to_2": "Peak absorbance ~2×10⁻⁴. The v₃ = 1 population reaches 7.2×10⁻⁴ (1 in 1400), "
                          "enough to produce measurable hot band absorption.",
            "nu3_2_to_3": "Peak absorbance ~2×10⁻⁷. Weak but nonzero — the v₃ = 2 population is "
                          "5.1×10⁻⁷ at 600 K.",
            "nu3_3_to_4": "Peak absorbance ~10⁻¹⁰. Extremely weak. The v₃ = 3 population at 600 K "
                          "is still negligibly small.",
        },
    ),
    (
        "combined_exomol_i1_absorbance_T1000K_P3Torr_x0p008_L100cm_step0p001",
        "Combined ExoMol I1 + HITRAN — 1000 K",
        "T = 1000 K, P = 3 Torr, x = 0.008, L = 100 cm, Δν = 0.001 cm⁻¹",
        {
            "nu3_0_to_1": "Fundamental continues to weaken (peak ~0.06). The Q(T)/Q(296) ratio "
                          "significantly dilutes all transitions.",
            "nu3_1_to_2": "Peak absorbance ~6×10⁻⁴. This is near the peak for the 1→2 band — "
                          "1.3% of molecules now occupy v₃ = 1, but the partition function dilution "
                          "will soon dominate.",
            "nu3_2_to_3": "Peak absorbance ~10⁻⁵. The v₃ = 2 population reaches 1.7×10⁻⁴, "
                          "making this band clearly detectable.",
            "nu3_3_to_4": "Peak absorbance ~9×10⁻⁸. Finally nonzero and measurable. "
                          "The v₃ = 3 population at 1000 K is ~2.2×10⁻⁶.",
        },
    ),
    (
        "combined_exomol_i1_absorbance_T2500K_P3Torr_x0p008_L100cm_step0p001",
        "Combined ExoMol I1 + HITRAN — 2500 K",
        "T = 2500 K, P = 3 Torr, x = 0.008, L = 100 cm, Δν = 0.001 cm⁻¹",
        {
            "nu3_0_to_1": "Fundamental is now very weak (peak ~3×10⁻⁴). The partition function at "
                          "2500 K is orders of magnitude larger than at 296 K, heavily diluting all transitions.",
            "nu3_1_to_2": "Peak absorbance ~7×10⁻⁵. The 1→2 band has passed its peak and is declining — "
                          "partition function dilution now outweighs the Boltzmann population gain.",
            "nu3_2_to_3": "Peak absorbance ~9×10⁻⁶. Still growing — the v₃ = 2 Boltzmann factor "
                          "has a steeper exponential rise (larger E″) so its peak temperature is higher.",
            "nu3_3_to_4": "Peak absorbance ~8×10⁻⁷. Still growing at 2500 K. The v₃ = 3 band "
                          "would peak at even higher temperatures.",
        },
    ),
    (
        "combined_exomol_i1_absorbance_T2500K_P3Torr_x1p0_L1000cm_step0p001",
        "Combined ExoMol I1 + HITRAN — 2500 K (pure CH₄, L = 1000 cm)",
        "T = 2500 K, P = 3 Torr, x = 1.0, L = 1000 cm, Δν = 0.001 cm⁻¹",
        {
            "nu3_0_to_1": "Pure CH₄ with 10× longer path amplifies the fundamental by ~125×. "
                          "Even at 2500 K the fundamental is clearly visible under these conditions.",
            "nu3_1_to_2": "The amplified conditions (x=1.0, L=1000 cm) make the 1→2 hot band "
                          "clearly visible. Compare with the dilute case at the same temperature.",
            "nu3_2_to_3": "The 2→3 band becomes clearly resolved with pure CH₄ and long path length. "
                          "Individual J-pair structure is now visible.",
            "nu3_3_to_4": "Even the 3→4 band shows resolved structure under these amplified conditions. "
                          "This demonstrates that the transitions exist but require extreme sensitivity to detect.",
        },
    ),
    (
        "exomol_direct_nu3_absorbance_T600K_P3Torr_x0p008_L100cm_step0p001",
        "ExoMol Direct — 600 K",
        "T = 600 K, P = 3 Torr, x = 0.008, L = 100 cm, Δν = 0.001 cm⁻¹\n"
        "Rendered directly from ExoMol MM database (γ₀ = 0.0488 cm⁻¹/bar, n = 0.4)",
        {
            "nu3_0_to_1": "ExoMol direct rendering at 600 K. Line intensities computed at the target "
                          "temperature using LTE formula, not scaled from 296 K reference. "
                          "Broadening uses ExoMol default parameters from .def file.",
            "nu3_1_to_2": "Compare with the Combined I1 result at 600 K. Differences arise from "
                          "the broadening model: ExoMol uses a single γ₀ for all lines, while "
                          "HAPI uses per-line broadening parameters.",
            "nu3_2_to_3": "ExoMol direct rendering of the 2→3 hot band. The line intensities are "
                          "computed at 600 K using the ExoMol partition function (covers up to 5000 K).",
            "nu3_3_to_4": "ExoMol direct rendering of the 3→4 hot band at 600 K. "
                          "This workflow can reach temperatures beyond HAPI's 2500 K limit.",
        },
    ),
]


def add_title_slide(prs: Presentation) -> None:
    layout = prs.slide_layouts[1]  # Title Slide
    slide = prs.slides.add_slide(layout)
    title = slide.placeholders[0]
    subtitle = slide.placeholders[1]
    title.text = "CH₄ ν₃ Absorbance Progressions\nTemperature Dependence"
    subtitle.text = (
        "Combined ExoMol MM I1 + HITRAN and ExoMol Direct workflows\n"
        "Progressions: 0→1 (fundamental), 1→2, 2→3, 3→4 (hot bands)\n"
        "Conditions: P = 3 Torr, x = 0.008, L = 100 cm, Δν = 0.001 cm⁻¹\n"
        "Temperatures: 300 K, 600 K, 1000 K, 2500 K"
    )


def add_overview_slide(prs: Presentation) -> None:
    layout = prs.slide_layouts[0]  # Title and Content
    slide = prs.slides.add_slide(layout)
    title = slide.placeholders[0]
    title.text = "Why hot bands vanish at low temperature"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    from pptx.dml.color import RGBColor

    tf.text = "Boltzmann population of v₃ = n relative to ground state"

    # (text, font_size_pt, bold, monospace)
    lines = [
        ("", 10, False, False),
        ("N(v₃ = n) / N(v₃ = 0) = exp(−n · 4343 / T)", 16, True, True),
        ("θ_vib = c₂ · ν₃ = 1.4388 × 3019 ≈ 4343 K", 14, False, True),
        ("", 8, False, False),
        ("   T (K)     v₃=1 / v₃=0     v₃=2 / v₃=0", 13, True, True),
        ("   ──────    ────────────     ────────────", 13, False, True),
        ("     296      4.2 × 10⁻⁷      1.8 × 10⁻¹³", 13, False, True),
        ("     600      7.2 × 10⁻⁴      5.1 × 10⁻⁷", 13, False, True),
        ("   1000      1.3 × 10⁻²      1.7 × 10⁻⁴", 13, False, True),
        ("   2000      1.1 × 10⁻¹      1.3 × 10⁻²", 13, False, True),
        ("   2500      1.8 × 10⁻¹      3.1 × 10⁻²", 13, False, True),
        ("", 8, False, False),
        ("At 1000 K, only 1.3% of molecules occupy v₃ = 1.", 14, False, False),
        ("The ν₃ mode (θ_vib = 4343 K) remains largely frozen out below ~2000 K.", 14, False, False),
        ("", 8, False, False),
        ("S(T) = S(T_ref) · Q(T_ref)/Q(T) · exp(−c₂E″/T)/exp(−c₂E″/T_ref) · [1−exp(−c₂ν₀/T)]/[1−exp(−c₂ν₀/T_ref)]", 12, False, True),
        ("Three competing effects: Boltzmann population ↑, partition function dilution ↓, stimulated emission (~1)", 12, False, False),
    ]
    for text, size, bold, mono in lines:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        if mono:
            p.font.name = "Consolas"


def add_figure_slide(
    prs: Presentation,
    *,
    png_path: Path,
    title_text: str,
    analysis_text: str,
    case_text: str,
) -> None:
    layout = prs.slide_layouts[0]  # Title and Content
    slide = prs.slides.add_slide(layout)

    title = slide.placeholders[0]
    title.text = title_text

    from pptx.dml.color import RGBColor
    from PIL import Image

    # Remove the content placeholder — we'll add image and text manually
    content_ph = slide.placeholders[1]
    sp = content_ph._element
    sp.getparent().remove(sp)

    # Compute image size to fit within available area
    img_obj = Image.open(png_path)
    aspect = img_obj.size[0] / img_obj.size[1]
    img_height = Inches(5.3)
    img_width = int(img_height * aspect)
    img_left = Inches(0.3)
    img_top = Inches(1.0)
    slide.shapes.add_picture(str(png_path), img_left, img_top, width=img_width, height=img_height)

    # Add analysis text box to the right of the image
    text_left = Inches(0.3 + img_width / 914400 + 0.3)
    text_width = Inches(13.33 - 0.3 - img_width / 914400 - 0.6)
    txBox = slide.shapes.add_textbox(text_left, Inches(1.2), text_width, Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = analysis_text
    p.font.size = Pt(12)

    # Add case parameters at the bottom
    footer_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.6), Inches(12.5), Inches(0.7))
    ftf = footer_box.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.text = case_text
    fp.font.size = Pt(9)
    fp.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    fp.alignment = PP_ALIGN.LEFT


def add_summary_slide(prs: Presentation) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    title = slide.placeholders[0]
    title.text = "Summary: Peak absorbance vs temperature"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    lines = [
        "Peak absorbance (max over all J pairs) at P=3 Torr, x=0.008, L=100 cm:",
        "",
        "Progression     300 K        600 K        1000 K       2500 K",
        "─────────────────────────────────────────────────────────────",
        "ν₃ 0→1          1.61         2.9×10⁻¹     5.9×10⁻²     3.0×10⁻⁴",
        "ν₃ 1→2          1.0×10⁻⁶     2.0×10⁻⁴     6.5×10⁻⁴     6.6×10⁻⁵",
        "ν₃ 2→3          7.2×10⁻¹³    2.4×10⁻⁷     1.2×10⁻⁵     9.4×10⁻⁶",
        "ν₃ 3→4          0.0          1.0×10⁻¹⁰    8.6×10⁻⁸     8.0×10⁻⁷",
        "",
        "Key observations:",
        "• Fundamental weakens with T due to partition function dilution (Q(T_ref)/Q(T))",
        "• 1→2 hot band peaks near 1000 K then declines — competition between",
        "  Boltzmann population growth and partition function dilution",
        "• Higher hot bands (2→3, 3→4) peak at progressively higher temperatures",
        "• ExoMol PF covers up to 5000 K; HAPI TIPS2025 limited to 2500 K",
        "• HITRAN lacks CH₄ ν₃ hot band data — only ExoMol contributes to 1→2, 2→3, 3→4",
    ]
    tf.text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
    for p in tf.paragraphs:
        p.font.size = Pt(13)


def main() -> int:
    prs = Presentation(str(TEMPLATE_PATH))

    # Remove existing template slides
    for _ in range(len(prs.slides)):
        rId = prs.slides._sldIdLst[-1].attrib[
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        ]
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

    add_title_slide(prs)
    add_overview_slide(prs)

    for folder_name, case_label, case_desc, analyses in CASES:
        folder = ARTIFACTS_DIR / folder_name
        if not folder.exists():
            print(f"SKIP (not found): {folder}")
            continue

        for prog in PROGRESSIONS:
            png = folder / f"{prog}_absorbance.png"
            if not png.exists():
                print(f"SKIP (no PNG): {png}")
                continue

            prog_label = PROGRESSION_LABELS[prog]
            analysis = analyses.get(prog, "")
            add_figure_slide(
                prs,
                png_path=png,
                title_text=f"{prog_label} — {case_label}",
                analysis_text=analysis,
                case_text=case_desc,
            )
            print(f"Added: {prog_label} — {case_label}")

    add_summary_slide(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"\nSaved: {OUTPUT_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
